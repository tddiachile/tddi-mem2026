"""
Exportar presentación TDDI a PowerPoint (.pptx) usando Selenium + Chrome.

Uso:
    pip install selenium python-pptx Pillow
    python export_pptx.py

Opciones:
    python export_pptx.py --output mi_presentacion.pptx
    python export_pptx.py --lang en
    python export_pptx.py --slides 12
"""

import argparse
import http.server
import io
import json
import os
import socket
import sys
import threading
import time


def find_free_port(start=8080):
    for port in range(start, start + 100):
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            if s.connect_ex(("127.0.0.1", port)) != 0:
                return port
    raise RuntimeError("No se encontró un puerto libre")


def start_server(directory, port):
    os.chdir(directory)
    handler = http.server.SimpleHTTPRequestHandler
    httpd = http.server.HTTPServer(("127.0.0.1", port), handler)
    thread = threading.Thread(target=httpd.serve_forever, daemon=True)
    thread.start()
    return httpd


def export_pptx(url, output_path, lang="es", total_slides=12):
    from pptx import Presentation  # noqa: F401 — validar antes de capturar
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options
    from selenium.webdriver.common.by import By
    from selenium.webdriver.common.keys import Keys

    print(f"  Abriendo Chrome headless...")

    opts = Options()
    opts.add_argument("--headless=new")
    opts.add_argument("--window-size=1920,1080")
    opts.add_argument("--disable-gpu")
    opts.add_argument("--no-sandbox")
    opts.add_argument("--ignore-certificate-errors")
    opts.add_argument("--disable-dev-shm-usage")
    opts.add_argument("--force-device-scale-factor=0.9")

    driver = webdriver.Chrome(options=opts)

    try:
        print(f"  Cargando {url}")
        driver.get(url)
        time.sleep(4)

        if lang == "en":
            try:
                btn = driver.find_element(By.XPATH, "//button[text()='EN']")
                btn.click()
                time.sleep(1)
            except Exception:
                pass

        # Cerrar el menu lateral
        driver.execute_script("""
            document.dispatchEvent(new MouseEvent('mousedown', {bubbles: true, clientX: 800, clientY: 400}));
        """)
        time.sleep(0.5)

        # Pre-renderizar todas las slides
        print(f"  Pre-renderizando {total_slides} diapositivas...")
        body = driver.find_element(By.TAG_NAME, "body")
        for i in range(total_slides - 1):
            body.send_keys(Keys.ARROW_RIGHT)
            time.sleep(1.2)

        # Volver al slide 0
        driver.execute_script("""
            const dots = document.querySelectorAll('.slide-dot');
            if (dots.length > 0) dots[0].click();
        """)
        time.sleep(1.5)

        # Cargar títulos desde locales.json
        locales_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "locales.json")
        with open(locales_path, "r", encoding="utf-8") as f:
            locales = json.load(f)
        slide_titles = locales.get(lang, {}).get("segments", [])
        # Rellenar si faltan
        while len(slide_titles) < total_slides:
            slide_titles.append(f"Slide {len(slide_titles) + 1}")

        # Ocultar header y footer
        driver.execute_script("""
            const header = document.querySelector('.site-header');
            const footer = document.querySelector('.slide-footer');
            if (header) { header.style.height = '0'; header.style.overflow = 'hidden'; header.style.padding = '0'; header.style.border = 'none'; }
            if (footer) { footer.style.height = '0'; footer.style.overflow = 'hidden'; footer.style.padding = '0'; footer.style.border = 'none'; }
        """)
        time.sleep(0.5)

        # Capturar cada slide
        screenshots = []
        print("  Capturando diapositivas...")

        for i in range(total_slides):
            driver.execute_script(f"""
                const dots = document.querySelectorAll('.slide-dot');
                if (dots[{i}]) dots[{i}].click();
            """)
            time.sleep(1.5)

            png = driver.get_screenshot_as_png()
            screenshots.append(png)
            print(f"    Slide {i + 1}/{total_slides} OK")

        # Ensamblar PPTX
        print("  Ensamblando PowerPoint...")
        create_pptx(screenshots, slide_titles, output_path)

        print(f"\n  [OK] PPTX generado: {output_path}")
        print(f"     {len(screenshots)} diapositivas - widescreen 16:9")

    finally:
        driver.quit()


def create_pptx(png_list, titles, output_path):
    from pptx import Presentation
    from pptx.util import Emu, Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from PIL import Image

    prs = Presentation()

    # Formato widescreen 16:9
    prs.slide_width = Emu(12192000)   # 13.333 pulgadas
    prs.slide_height = Emu(6858000)   # 7.5 pulgadas

    blank_layout = prs.slide_layouts[6]  # Layout en blanco
    title_height = Inches(0.6)  # Altura reservada para el título

    for idx, png_bytes in enumerate(png_list):
        slide = prs.slides.add_slide(blank_layout)
        title_text = titles[idx] if idx < len(titles) else f"Slide {idx + 1}"

        # Agregar título en la parte superior
        txBox = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.1), prs.slide_width - Inches(0.8), title_height
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x2A, 0x29, 0x5C)  # sdx-dark
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.LEFT

        # Redimensionar imagen debajo del título
        img = Image.open(io.BytesIO(png_bytes))
        img_width, img_height = img.size

        slide_w = prs.slide_width
        available_h = prs.slide_height - title_height - Inches(0.15)

        # Escalar manteniendo aspect ratio
        scale_w = slide_w / img_width
        scale_h = available_h / img_height
        scale = min(scale_w, scale_h)

        final_w = int(img_width * scale)
        final_h = int(img_height * scale)
        left = (slide_w - final_w) // 2
        top = title_height + Inches(0.15)

        img_stream = io.BytesIO(png_bytes)
        slide.shapes.add_picture(img_stream, left, top, final_w, final_h)

    if not png_list:
        print("  [WARN] No se capturaron imágenes")
        return

    prs.save(output_path)


def main():
    parser = argparse.ArgumentParser(description="Exportar presentación TDDI a PowerPoint")
    parser.add_argument("--output", "-o", default="presentacion_tddi.pptx", help="Archivo PPTX de salida")
    parser.add_argument("--port", "-p", type=int, default=0, help="Puerto del servidor (0=auto)")
    parser.add_argument("--lang", "-l", choices=["es", "en"], default="es", help="Idioma")
    parser.add_argument("--slides", "-s", type=int, default=12, help="Número de diapositivas")
    parser.add_argument("--url", "-u", default=None, help="URL directa (omite servidor local)")
    args = parser.parse_args()

    project_dir = os.path.dirname(os.path.abspath(__file__))

    print("\n  Exportador de Presentacion TDDI -> PowerPoint")
    print("  ------------------------------------------------\n")

    httpd = None
    if args.url:
        url = args.url
    else:
        port = args.port or find_free_port()
        print(f"  Servidor local en puerto {port}...")
        httpd = start_server(project_dir, port)
        url = f"http://127.0.0.1:{port}/index.html"

    try:
        output = os.path.join(project_dir, args.output)
        export_pptx(url, output, lang=args.lang, total_slides=args.slides)
    except ImportError as e:
        print(f"\n  [WARN] Dependencia faltante: {e}")
        print("    pip install selenium python-pptx Pillow\n")
        sys.exit(1)
    except Exception as e:
        print(f"\n  [ERROR] {e}")
        sys.exit(1)
    finally:
        if httpd:
            httpd.shutdown()


if __name__ == "__main__":
    main()
